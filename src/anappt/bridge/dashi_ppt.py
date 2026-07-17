"""Dashi-PPT bridge layer for AnaPPTAgent.

Converts analysis report markdown into HTML-based slide presentations.
Generates self-contained HTML files that can be opened in a browser
and exported to PPTX/PDF.
"""

from __future__ import annotations

import re
from pathlib import Path

from anappt.i18n import t

# Available presentation themes
_THEMES: dict[str, str] = {
    "default": "Default - Clean white background with blue accents",
    "dark": "Dark - Dark background with light text",
    "corporate": "Corporate - Professional blue theme",
    "minimal": "Minimal - Simple black and white",
    "vibrant": "Vibrant - Colorful with gradients",
}

# CSS for each theme
_THEME_CSS: dict[str, str] = {
    "default": """
        background: #ffffff; color: #333;
        --accent: #2563eb; --bg-alt: #f0f4ff;
    """,
    "dark": """
        background: #1a1a2e; color: #e0e0e0;
        --accent: #818cf8; --bg-alt: #16213e;
    """,
    "corporate": """
        background: #ffffff; color: #1e3a5f;
        --accent: #1e40af; --bg-alt: #eff6ff;
    """,
    "minimal": """
        background: #ffffff; color: #000000;
        --accent: #555555; --bg-alt: #f5f5f5;
    """,
    "vibrant": """
        background: linear-gradient(135deg, #667eea, #764ba2);
        color: #ffffff; --accent: #ffd700; --bg-alt: rgba(255,255,255,0.1);
    """,
}


class SlideContent:
    """Represents a single slide's content.

    Attributes:
        title: Slide title.
        bullets: List of bullet point strings.
        content: Optional raw markdown content block.
        image_path: Optional path to an image.
        layout: Slide layout type ('title', 'content', 'image', 'section').
    """

    def __init__(
        self,
        title: str = "",
        bullets: list[str] | None = None,
        content: str = "",
        image_path: str = "",
        layout: str = "content",
    ) -> None:
        """Initialize slide content.

        Args:
            title: Slide title.
            bullets: List of bullet points.
            content: Raw content block.
            image_path: Path to an image file.
            layout: Layout type.
        """
        self.title: str = title
        self.bullets: list[str] = bullets or []
        self.content: str = content
        self.image_path: str = image_path
        self.layout: str = layout


class DashiPPTBridge:
    """Bridge to generate HTML-based slide presentations.

    Converts markdown report content into a structured HTML presentation
    with theme support. The generated HTML is self-contained and can be
    opened directly in a browser.

    Attributes:
        output_dir: Directory for generated PPT files.
        theme: Selected theme name.
    """

    def __init__(
        self,
        output_dir: str | Path | None = None,
        theme: str = "default",
    ) -> None:
        """Initialize the bridge.

        Args:
            output_dir: Output directory for generated files.
                        Defaults to current working directory.
            theme: Theme name from available themes.
        """
        self.output_dir: Path = Path(output_dir) if output_dir else Path.cwd()
        self.theme: str = theme if theme in _THEMES else "default"

    @staticmethod
    def list_themes() -> dict[str, str]:
        """Return available themes.

        Returns:
            Dictionary mapping theme name to description.
        """
        return dict(_THEMES)

    @staticmethod
    def validate_markdown(content: str) -> bool:
        """Validate that markdown content can be converted to slides.

        Checks that the content has at least one heading (H1 or H2)
        which will become a slide title.

        Args:
            content: Markdown content string.

        Returns:
            True if the content is valid for slide generation.
        """
        if not content or not content.strip():
            return False
        # Must have at least one heading
        has_heading = bool(re.search(r"^#{1,3}\s+", content, re.MULTILINE))
        return has_heading

    def parse_markdown_to_slides(self, markdown: str) -> list[SlideContent]:
        """Parse markdown content into a list of slides.

        Each H1 or H2 heading starts a new slide. Content under a heading
        becomes the slide body. Bullet points (lines starting with - or *)
        are extracted as bullets.

        Args:
            markdown: Markdown content string.

        Returns:
            List of SlideContent objects.
        """
        slides: list[SlideContent] = []
        lines = markdown.split("\n")

        current_title = ""
        current_bullets: list[str] = []
        current_content_lines: list[str] = []

        def flush_slide() -> None:
            nonlocal current_title, current_bullets, current_content_lines
            if current_title or current_bullets or current_content_lines:
                layout = "title" if not current_bullets and not current_content_lines else "content"
                slides.append(
                    SlideContent(
                        title=current_title,
                        bullets=list(current_bullets),
                        content="\n".join(current_content_lines),
                        layout=layout,
                    )
                )
            current_title = ""
            current_bullets = []
            current_content_lines = []

        for line in lines:
            stripped = line.strip()

            # H1 or H2 starts a new slide
            h1_match = re.match(r"^#\s+(.+)$", stripped)
            h2_match = re.match(r"^##\s+(.+)$", stripped)

            if h1_match:
                flush_slide()
                current_title = h1_match.group(1).strip()
            elif h2_match:
                flush_slide()
                current_title = h2_match.group(1).strip()
            elif stripped.startswith("- ") or stripped.startswith("* "):
                bullet = stripped[2:].strip()
                current_bullets.append(bullet)
            elif stripped:
                current_content_lines.append(stripped)

        flush_slide()
        return slides

    def generate_html(
        self,
        slides: list[SlideContent],
        title: str = "Presentation",
    ) -> str:
        """Generate a self-contained HTML presentation from slides.

        Args:
            slides: List of SlideContent objects.
            title: Presentation title (for the HTML page title).

        Returns:
            Complete HTML string.
        """
        theme_css = _THEME_CSS.get(self.theme, _THEMES["default"])

        slides_html: list[str] = []
        for i, slide in enumerate(slides):
            slide_class = "slide" + (" active" if i == 0 else "")
            inner: list[str] = []

            if slide.title:
                inner.append(f'<h1 class="slide-title">{slide.title}</h1>')

            if slide.bullets:
                bullets_html = "\n".join(f"<li>{b}</li>" for b in slide.bullets)
                inner.append(f'<ul class="slide-bullets">{bullets_html}</ul>')

            if slide.content:
                inner.append(f'<div class="slide-content">{slide.content}</div>')

            inner_str = "\n".join(inner)
            slides_html.append(
                f'<section class="{slide_class}" data-index="{i}">\n{inner_str}\n</section>'
            )

        slides_joined = "\n".join(slides_html)

        return f"""<!DOCTYPE html>
<html lang="zh">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>{title}</title>
    <style>
        * {{ margin: 0; padding: 0; box-sizing: border-box; }}
        body {{
            {theme_css}
            font-family: -apple-system, BlinkMacSystemFont, 'Segoe UI', Roboto, sans-serif;
            min-height: 100vh;
            overflow: hidden;
        }}
        .presentation {{
            position: relative;
            width: 100vw;
            height: 100vh;
        }}
        .slide {{
            display: none;
            position: absolute;
            top: 0; left: 0;
            width: 100%; height: 100%;
            padding: 60px 80px;
            flex-direction: column;
            justify-content: center;
        }}
        .slide.active {{ display: flex; }}
        .slide-title {{
            font-size: 2.5em;
            margin-bottom: 30px;
            color: var(--accent);
            border-bottom: 3px solid var(--accent);
            padding-bottom: 10px;
        }}
        .slide-bullets {{
            font-size: 1.4em;
            line-height: 2;
            list-style: none;
        }}
        .slide-bullets li {{
            padding: 8px 0;
            padding-left: 30px;
            position: relative;
        }}
        .slide-bullets li::before {{
            content: '\\25B8';
            position: absolute;
            left: 0;
            color: var(--accent);
        }}
        .slide-content {{
            font-size: 1.2em;
            line-height: 1.8;
            margin-top: 20px;
            white-space: pre-wrap;
        }}
        .nav {{
            position: fixed;
            bottom: 20px; right: 30px;
            display: flex;
            gap: 10px;
            z-index: 100;
        }}
        .nav button {{
            padding: 8px 20px;
            font-size: 1.1em;
            border: 2px solid var(--accent);
            background: transparent;
            color: var(--accent);
            border-radius: 8px;
            cursor: pointer;
        }}
        .nav button:hover {{
            background: var(--accent);
            color: var(--bg-alt);
        }}
        .counter {{
            position: fixed;
            bottom: 25px; left: 30px;
            font-size: 1.1em;
            color: var(--accent);
        }}
    </style>
</head>
<body>
    <div class="presentation">
        {slides_joined}
    </div>
    <div class="counter"><span id="current">1</span> / <span id="total">{len(slides)}</span></div>
    <div class="nav">
        <button onclick="prev()">&#8592; Prev</button>
        <button onclick="next()">Next &#8594;</button>
    </div>
    <script>
        let current = 0;
        const slides = document.querySelectorAll('.slide');
        const total = slides.length;

        function show(index) {{
            slides.forEach((s, i) => {{
                s.classList.toggle('active', i === index);
            }});
            document.getElementById('current').textContent = index + 1;
        }}

        function next() {{
            if (current < total - 1) {{
                current++;
                show(current);
            }}
        }}

        function prev() {{
            if (current > 0) {{
                current--;
                show(current);
            }}
        }}

        document.addEventListener('keydown', (e) => {{
            if (e.key === 'ArrowRight' || e.key === ' ') next();
            if (e.key === 'ArrowLeft') prev();
        }});
    </script>
</body>
</html>"""

    def generate_ppt(
        self,
        markdown_content: str,
        theme: str | None = None,
        title: str = "Analysis Report",
        filename: str = "presentation.html",
    ) -> Path:
        """Generate a complete HTML presentation from markdown.

        Args:
            markdown_content: Markdown report content.
            theme: Theme name (uses self.theme if None).
            title: Presentation title.
            filename: Output filename.

        Returns:
            Path to the generated HTML file.

        Raises:
            ValueError: If markdown content is invalid.
        """
        if not self.validate_markdown(markdown_content):
            raise ValueError(t("bridge.markdown_invalid"))

        if theme is not None:
            self.theme = theme if theme in _THEMES else "default"

        slides = self.parse_markdown_to_slides(markdown_content)
        html = self.generate_html(slides, title=title)

        self.output_dir.mkdir(parents=True, exist_ok=True)
        output_path = self.output_dir / filename
        output_path.write_text(html, encoding="utf-8")
        return output_path

    def generate_pptx(
        self,
        markdown_content: str,
        theme: str | None = None,
        title: str = "Analysis Report",
        filename: str = "presentation.pptx",
    ) -> Path:
        """Generate a PPTX file from markdown content.

        Uses python-pptx if available; otherwise falls back to HTML.

        Args:
            markdown_content: Markdown report content.
            theme: Theme name.
            title: Presentation title.
            filename: Output filename.

        Returns:
            Path to the generated file.
        """
        if not self.validate_markdown(markdown_content):
            raise ValueError(t("bridge.markdown_invalid"))

        if theme is not None:
            self.theme = theme if theme in _THEMES else "default"

        slides = self.parse_markdown_to_slides(markdown_content)

        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
        except ImportError:
            # Fall back to HTML if python-pptx not installed
            return self.generate_ppt(
                markdown_content, theme=theme, title=title, filename="presentation.html"
            )

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        for slide_content in slides:
            slide_layout = prs.slide_layouts[1]  # Title and Content
            slide = prs.slides.add_slide(slide_layout)

            if slide_content.title:
                slide.shapes.title.text = slide_content.title

            if slide_content.bullets:
                body_shape = slide.shapes.placeholders[1]
                tf = body_shape.text_frame
                for i, bullet in enumerate(slide_content.bullets):
                    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                    p.text = bullet
                    p.font.size = Pt(18)

        self.output_dir.mkdir(parents=True, exist_ok=True)
        output_path = self.output_dir / filename
        prs.save(str(output_path))
        return output_path
